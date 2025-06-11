from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 default dimensions
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullet_slide(title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for idx, line in enumerate(bullets):
        p = tf.add_paragraph() if idx else tf.paragraphs[0]
        p.text = line
        p.level = 0
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

# Slide 1 - Title
add_title_slide("JP Morgan Chase Financial AI Agent",
                "Hackathon Spring 2025  •  Team …  |  April 27 2025")

# Slide 2 - Agenda
add_bullet_slide("Agenda", [
    "Problem & Goals",
    "Solution Overview",
    "Architecture & Key Components",
    "Demo Highlights",
    "Impact & KPI Alignment",
    "Road‑map & Q&A"
])

# Slide 3 - Problem & Goals
add_bullet_slide("Business Problem", [
    "No digital financial‑advisor experience",
    "Customers cannot upload / analyse 10‑K, 10‑Q reports",
    "High call‑center load, low self‑service rate"
], "Rubric 15 % : problem intro, KPIs – reduce support calls 25 %, lift NPS +10")

# Slide 4 - Goals ↔ Hackathon rubric
add_bullet_slide("Hackathon Goals Tie‑In", [
    "Automation – chatbot answers routine questions 24×7",
    "Personalisation – verifies customer & serves account data",
    "Operational Efficiency – PDF insight lowers manual effort",
    "Presentation: clarity, depth, actionable recommendations"
])

# Slide 5 - Solution Overview
add_bullet_slide("Solution Overview", [
    "Gradio front‑end + OpenAI GPT‑3.5 reasoning back‑end",
    "FSM verifies existing customers (phone → ZIP)",
    "Prospect funnel captures name • phone • email",
    "Single‑PDF upload → real‑time Q&A on 10‑K / 10‑Q",
    "Balance, recent transactions, savings offers"
])

# Slide 6 - Architecture
add_bullet_slide("High‑Level Architecture", [
    "Front‑end: Python Gradio UI (chat + file upload)",
    "Logic: Finite‑state machine for conversation flow",
    "Data: SQLite customers.db & pdfplumber extraction",
    "AI: OpenAI Chat Completions API",
    "Hosting: single Python script, lightweight deploy"
], "Diagram on screen during live demo")

# Slide 7 - Code Highlights (quality / uniqueness 35 %) 
add_bullet_slide("Key Code Highlights", [
    "≈200 LOC, fully self‑contained",
    "Dynamic follow‑ups: different prompts for new vs existing",
    "Fuzzy intent detection (regex) for ‘existing / new’",
    "Error‑handling & exit confirmation for polished UX",
    "8 k‑char context window – cost‑controlled PDF analysis"
])

# Slide 8 - Demo Flow (2‑min narrated video)
add_bullet_slide("Demo Snapshot", [
    "Existing customer → balance & transactions",
    "Upload 10‑K PDF → ask “operating expenses 2024?”",
    "New user → lead capture funnel"
], "Narrate 3 short screencaps – 30 s each")

# Slide 9 - Impact & KPIs
add_bullet_slide("Impact & KPI Alignment", [
    "↓ 25 % projected call‑center volume (automation)",
    "↑ 10 pt Net Promoter Score (personalised answers)",
    "↓ 70 % manual effort reading financial docs",
    "Scalable: same pattern for mortgages, wealth, etc."
])

# Slide 10 - Road‑map & Closing
add_bullet_slide("Next Steps & Q&A", [
    "Multi‑PDF support & semantic search vector DB",
    "Add voice interface & Salesforce lead integration",
    "Deploy on Azure App Service with CI/CD",
    "Questions?"
], "End with thank‑you and contact e‑mails")

# Save presentation
file_path = "/mnt/data/JP_Morgan_FinAI_Agent.pptx"
prs.save(file_path)

file_path
